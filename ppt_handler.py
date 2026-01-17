import os
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData, XyChartData
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.series import SeriesCollection
from typing import Dict, Any, List, Tuple
import pandas as pd
import numpy as np
from io import BytesIO
from rbc_branding import RBCBranding

class PowerPointHandler:
    def __init__(self, template_path: str = None):
        self.template_path = template_path
        self.presentation = None
        self.chart_style_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'area': XL_CHART_TYPE.AREA,
            'histogram': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'stacked_bar': XL_CHART_TYPE.COLUMN_STACKED,
            'stacked_area': XL_CHART_TYPE.AREA
        }
        
        # Set matplotlib style with RBC branding
        RBCBranding.setup_matplotlib_style()
    
    def load_template(self, template_path: str = None):
        """Load PowerPoint template or create new presentation."""
        if template_path:
            self.template_path = template_path
        
        if self.template_path and os.path.exists(self.template_path):
            self.presentation = Presentation(self.template_path)
        else:
            self.presentation = Presentation()
    
    def create_chart_image(self, data: pd.DataFrame, chart_config: Dict[str, Any], 
                          width: int = 10, height: int = 6) -> BytesIO:
        """Create a chart image using matplotlib and return as BytesIO."""
        chart_type = chart_config['chart_type']
        
        fig, ax = plt.subplots(figsize=(width, height))
        fig.patch.set_facecolor('white')
        
        if chart_type == 'bar':
            self._create_bar_chart(ax, data, chart_config)
        elif chart_type == 'line':
            self._create_line_chart(ax, data, chart_config)
        elif chart_type == 'pie':
            self._create_pie_chart(ax, data, chart_config)
        elif chart_type == 'scatter':
            self._create_scatter_chart(ax, data, chart_config)
        elif chart_type == 'area':
            self._create_area_chart(ax, data, chart_config)
        elif chart_type == 'histogram':
            self._create_histogram(ax, data, chart_config)
        elif chart_type == 'stacked_bar':
            self._create_stacked_bar_chart(ax, data, chart_config)
        elif chart_type == 'stacked_area':
            self._create_stacked_area_chart(ax, data, chart_config)
        
        plt.tight_layout()
        
        # Save to BytesIO
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close()
        
        return img_buffer
    
    def _create_bar_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create bar chart with RBC branding."""
        x_col = config['x_column']
        y_col = config['y_column']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('bar')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'bar')
        
        # Create bars with RBC colors
        colors = RBCBranding.get_color_palette('categorical', len(data))
        bars = ax.bar(data[x_col], data[y_col], color=colors, 
                   edgecolor=style_config['edge_color'], linewidth=style_config['edge_width'])
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel(y_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title(f'{y_col} by {x_col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Rotate x-axis labels if needed
        if len(data[x_col].unique()) > 5:
            plt.xticks(rotation=45, ha='right', fontname=style_config['label_font'], 
                      fontsize=style_config['label_size'])
        
        # Add value labels on bars
        for i, bar in enumerate(bars):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}', ha='center', va='bottom',
                   fontname=style_config['label_font'], fontsize=style_config['label_size']-1)
    
    def _create_line_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create line chart with RBC branding."""
        x_col = config['x_column']
        y_col = config['y_column']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('line')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'line')
        
        # Create line with RBC colors
        line_color = RBCBranding.COLORS['primary_blue']
        ax.plot(data[x_col], data[y_col], marker='o', 
               color=line_color, linewidth=style_config['line_width'], 
               markersize=style_config['marker_size'],
               markerfacecolor=RBCBranding.COLORS['accent_yellow'],
               markeredgecolor=line_color, markeredgewidth=1.5)
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel(y_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title(f'{y_col} over {x_col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Format x-axis for dates
        if data[x_col].dtype == 'datetime64[ns]':
            fig = ax.figure
            fig.autofmt_xdate()
    
    def _create_pie_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create pie chart with RBC branding."""
        labels_col = config['labels_column']
        values_col = config['values_column']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('pie')
        
        # Group data by labels and sum values
        grouped_data = data.groupby(labels_col)[values_col].sum()
        
        # Create pie with RBC colors
        colors = RBCBranding.get_color_palette('categorical', len(grouped_data))
        wedges, texts, autotexts = ax.pie(grouped_data.values, labels=grouped_data.index,
                                          autopct='%1.1f%%', startangle=90, colors=colors,
                                          wedgeprops=style_config['wedge_props'])
        
        # Set title with RBC fonts
        ax.set_title(f'Distribution of {values_col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Equal aspect ratio ensures that pie is drawn as a circle
        ax.axis('equal')
    
    def _create_scatter_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create scatter plot with RBC branding."""
        x_col = config['x_column']
        y_col = config['y_column']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('scatter')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'scatter')
        
        # Create scatter with RBC colors
        scatter_color = RBCBranding.COLORS['primary_blue']
        ax.scatter(data[x_col], data[y_col], color=scatter_color, 
                  alpha=style_config['scatter_alpha'], s=80,
                  edgecolor=style_config['scatter_edge_color'], 
                  linewidth=style_config['scatter_edge_width'])
        
        # Add trend line with RBC accent color
        z = np.polyfit(data[x_col], data[y_col], 1)
        p = np.poly1d(z)
        ax.plot(data[x_col], p(data[x_col]), color=RBCBranding.COLORS['accent_yellow'], 
               linestyle='--', linewidth=2, alpha=0.8)
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel(y_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title(f'{y_col} vs {x_col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
    
    def _create_area_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create area chart with RBC branding."""
        x_col = config['x_column']
        y_col = config['y_column']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('area')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'area')
        
        # Create area with RBC colors
        area_color = RBCBranding.COLORS['primary_blue']
        ax.fill_between(data[x_col], data[y_col], color=area_color, 
                      alpha=style_config['area_alpha'])
        ax.plot(data[x_col], data[y_col], color=area_color, 
               linewidth=style_config['line_width'])
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel(y_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title(f'{y_col} over {x_col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Format x-axis for dates
        if data[x_col].dtype == 'datetime64[ns]':
            fig = ax.figure
            fig.autofmt_xdate()
    
    def _create_histogram(self, ax, data: pd.DataFrame, config: Dict):
        """Create histogram with RBC branding."""
        col = config['column']
        bins = config.get('bins', 20)
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('histogram')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'histogram')
        
        # Create histogram with RBC colors
        ax.hist(data[col].dropna(), bins=bins, color=style_config['hist_color'], 
               alpha=style_config['hist_alpha'], 
               edgecolor=style_config['hist_edge_color'])
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel('Frequency', fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title(f'Distribution of {col}', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
    
    def _create_stacked_bar_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create stacked bar chart with RBC branding."""
        x_col = config['x_column']
        y_cols = config['y_columns']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('bar')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'bar')
        
        # Get RBC colors for each series
        colors = RBCBranding.get_color_palette('categorical', len(y_cols))
        
        # Create stacked bars
        bottom = np.zeros(len(data))
        for i, y_col in enumerate(y_cols):
            ax.bar(data[x_col], data[y_col], bottom=bottom, color=colors[i], 
                   edgecolor=style_config['edge_color'], linewidth=style_config['edge_width'],
                   label=y_col)
            bottom += data[y_col]
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel('Value', fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title('Stacked Bar Chart', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Add legend
        legend = ax.legend(fontsize=style_config['label_size']-1)
        # Set legend font properties
        for text in legend.get_texts():
            text.set_fontname(style_config['label_font'])
            text.set_fontsize(style_config['label_size']-1)
        
        # Rotate x-axis labels if needed
        if len(data[x_col].unique()) > 5:
            plt.xticks(rotation=45, ha='right', fontname=style_config['label_font'], 
                      fontsize=style_config['label_size'])
        
        # Format y-axis as percentage
        ax.set_ylim(0, 100)
        ax.set_ylabel('Percentage (%)', fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
    
    def _create_stacked_area_chart(self, ax, data: pd.DataFrame, config: Dict):
        """Create 100% stacked area chart with RBC branding."""
        x_col = config['x_column']
        y_cols = config['y_columns']
        
        # Get RBC styling
        style_config = RBCBranding.get_chart_style_config('area')
        RBCBranding.apply_rbc_theme_to_axis(ax, 'area')
        
        # Get RBC colors for each series
        colors = RBCBranding.get_color_palette('sequential', len(y_cols))
        
        # Create stacked area chart (100% stacked)
        # Normalize data to 100%
        data_normalized = data.copy()
        for y_col in y_cols:
            data_normalized[y_col] = (data[y_col] / data[y_cols].sum(axis=1)) * 100
        
        # Create stacked areas
        bottom = np.zeros(len(data))
        for i, y_col in enumerate(y_cols):
            ax.fill_between(data[x_col], bottom, bottom + data_normalized[y_col], 
                         color=colors[i], alpha=style_config['area_alpha'], label=y_col)
            bottom += data_normalized[y_col]
        
        # Set labels and title with RBC fonts
        ax.set_xlabel(x_col, fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_ylabel('Percentage (%)', fontname=style_config['label_font'], 
                    fontsize=style_config['label_size'], color=style_config['label_color'])
        ax.set_title('100% Stacked Area Chart', fontname=style_config['title_font'], 
                   fontsize=style_config['title_size'], color=style_config['title_color'], weight='bold')
        
        # Add legend
        legend = ax.legend(fontsize=style_config['label_size']-1)
        # Set legend font properties
        for text in legend.get_texts():
            text.set_fontname(style_config['label_font'])
            text.set_fontsize(style_config['label_size']-1)
        
        # Format y-axis as percentage
        ax.set_ylim(0, 100)
        
        # Format x-axis for dates
        if data[x_col].dtype == 'datetime64[ns]':
            fig = ax.figure
            fig.autofmt_xdate()
    
    def add_slide_with_chart(self, data: pd.DataFrame, chart_config: Dict[str, Any], 
                           title: str = None, slide_layout: int = 6) -> int:
        """Add a new slide with chart to the presentation."""
        if not self.presentation:
            self.load_template()
        
        # Add new slide
        slide_layout = self.presentation.slide_layouts[slide_layout]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title with RBC styling
        if title:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = title
                # Apply RBC font styling
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = RBCBranding.FONTS['heading']
                        run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                        run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                        run.font.bold = True
            else:
                # Create a text box for title if no title placeholder
                from pptx.util import Inches
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title
                # Apply RBC font styling
                for paragraph in title_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = RBCBranding.FONTS['heading']
                        run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                        run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                        run.font.bold = True
        else:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = f"{chart_config['chart_type'].title()} Chart"
                # Apply RBC font styling
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = RBCBranding.FONTS['heading']
                        run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                        run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                        run.font.bold = True
            else:
                # Create a text box for title if no title placeholder
                from pptx.util import Inches
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = f"{chart_config['chart_type'].title()} Chart"
                # Apply RBC font styling
                for paragraph in title_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = RBCBranding.FONTS['heading']
                        run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                        run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                        run.font.bold = True
        
        # Create chart image
        chart_img = self.create_chart_image(data, chart_config)
        
        # Add chart image to slide
        slide.shapes.add_picture(chart_img, Inches(1), Inches(1.5), Inches(8), Inches(5))
        
        return len(self.presentation.slides) - 1
    
    def add_summary_slide(self, data_summary: Dict[str, Any], title: str = "Data Summary") -> int:
        """Add a summary slide with key insights."""
        if not self.presentation:
            self.load_template()
        
        slide_layout = self.presentation.slide_layouts[1]  # Title and content
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = title
            # Apply RBC font styling to title
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = RBCBranding.FONTS['heading']
                    run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                    run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                    run.font.bold = True
        else:
            # Create a text box for title if no title placeholder
            from pptx.util import Inches
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            # Apply RBC font styling
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = RBCBranding.FONTS['heading']
                    run.font.size = Pt(RBCBranding.FONT_SIZES['title'])
                    run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
                    run.font.bold = True
        
        # Add summary content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()
        
        # Add summary points with RBC styling
        p = tf.add_paragraph()
        p.text = f"Dataset Shape: {data_summary['shape'][0]} rows × {data_summary['shape'][1]} columns"
        p.level = 0
        for run in p.runs:
            run.font.name = RBCBranding.FONTS['body']
            run.font.size = Pt(RBCBranding.FONT_SIZES['body'])
            run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
        
        p = tf.add_paragraph()
        p.text = f"Numeric Columns: {len(data_summary['numeric_columns'])}"
        p.level = 0
        for run in p.runs:
            run.font.name = RBCBranding.FONTS['body']
            run.font.size = Pt(RBCBranding.FONT_SIZES['body'])
            run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
        
        p = tf.add_paragraph()
        p.text = f"Categorical Columns: {len(data_summary['categorical_columns'])}"
        p.level = 0
        for run in p.runs:
            run.font.name = RBCBranding.FONTS['body']
            run.font.size = Pt(RBCBranding.FONT_SIZES['body'])
            run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
        
        if data_summary['missing_values']:
            total_missing = sum(data_summary['missing_values'].values())
            p = tf.add_paragraph()
            p.text = f"Missing Values: {total_missing}"
            p.level = 0
            for run in p.runs:
                run.font.name = RBCBranding.FONTS['body']
                run.font.size = Pt(RBCBranding.FONT_SIZES['body'])
                run.font.color.rgb = RGBColor(*RBCBranding.RGB_COLORS['dark_blue'])
        
        return len(self.presentation.slides) - 1
    
    def save_presentation(self, output_path: str):
        """Save the presentation to file."""
        if not self.presentation:
            raise ValueError("No presentation to save. Create slides first.")
        
        self.presentation.save(output_path)
    
    def create_presentation_from_data(self, data: pd.DataFrame, chart_configs: List[Dict[str, Any]], 
                                    title: str = "Data Visualization Report") -> str:
        """Create a complete presentation from data and chart configurations."""
        self.load_template()
        
        # Add title slide
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = title
        else:
            # Create a text box for title if no title placeholder
            from pptx.util import Inches
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
        
        subtitle_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if subtitle_shape is not None:
            subtitle_shape.text = f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Add summary slide
        from data_analyzer import DataAnalyzer
        analyzer = DataAnalyzer()
        data_summary = analyzer.analyze_data(data)
        self.add_summary_slide(data_summary)
        
        # Add chart slides
        for i, config in enumerate(chart_configs):
            chart_title = f"{config['chart_type'].title()} Chart {i+1}"
            self.add_slide_with_chart(data, config, chart_title)
        
        # Save presentation
        output_path = f"{title.replace(' ', '_')}.pptx"
        self.save_presentation(output_path)
        
        return output_path
